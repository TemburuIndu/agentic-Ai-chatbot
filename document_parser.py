def extract_pptx(path_or_url: str) -> str:
    from pptx import Presentation
    from PIL import Image
    import pytesseract
    from io import BytesIO
    import requests
    import os
    from concurrent.futures import ThreadPoolExecutor

    def ocr_image(img_bytes: bytes) -> str:
        img = Image.open(BytesIO(img_bytes))
        return pytesseract.image_to_string(img, lang='eng').strip()

    try:
        # Load presentation from URL or local path
        if path_or_url.lower().startswith(("http://", "https://")):
            prs = Presentation(BytesIO(requests.get(path_or_url).content))
        else:
            if not os.path.exists(path_or_url):
                raise FileNotFoundError(f"PPTX not found: {path_or_url}")
            prs = Presentation(path_or_url)

        all_text = []
        images_to_ocr = []

        # Single pass: extract text and collect images
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    all_text.append(shape.text.strip())
                elif shape.shape_type == 13 and hasattr(shape, "image"):
                    images_to_ocr.append(shape.image.blob)

        # Parallel OCR for images
        if images_to_ocr:
            with ThreadPoolExecutor() as executor:
                ocr_results = executor.map(ocr_image, images_to_ocr)
                all_text.extend([res for res in ocr_results if res])

        return "\n".join(all_text)

    except Exception as e:
        print(f"❌ Error processing PPTX {path_or_url}: {e}")
        return ""



# print(extract_pptx("https://hackrx.blob.core.windows.net/assets/Test%20/Test%20Case%20HackRx.pptx?sv=2023-01-03&spr=https&st=2025-08-04T18%3A36%3A56Z&se=2026-08-05T18%3A36%3A00Z&sr=b&sp=r&sig=v3zSJ%2FKW4RhXaNNVTU9KQbX%2Bmo5dDEIzwaBzXCOicJM%3D"))
# print(extract_pptx("/home/krishna/hackathon-agent/files/Test Case HackRx.pptx"))

def extract_xlsx(path):
    from markitdown import MarkItDown
    md = MarkItDown()
    result = md.convert(path)
    return result.markdown

# print(extract_xlsx("https://hackrx.blob.core.windows.net/assets/Test%20/Salary%20data.xlsx?sv=2023-01-03&spr=https&st=2025-08-04T18%3A46%3A54Z&se=2026-08-05T18%3A46%3A00Z&sr=b&sp=r&sig=sSoLGNgznoeLpZv%2FEe%2FEI1erhD0OQVoNJFDPtqfSdJQ%3D"))
# print(extract_xlsx("/home/krishna/hackathon-agent/files/Salary data.xlsx"))

def extract_image(path_or_url: str) -> str:
    from PIL import Image
    import pytesseract
    import requests
    from io import BytesIO
    import os

    """
    Extracts text from an image (local file or hosted URL) using OCR (Tesseract).
    Supports JPEG, PNG, BMP, TIFF, etc.
    No temporary local file is created for URLs.
    """

    try:
        if path_or_url.lower().startswith(("http://", "https://")):
            # Remote image
            response = requests.get(path_or_url, stream=True)
            response.raise_for_status()
            image = Image.open(BytesIO(response.content))
        else:
            # Local image
            if not os.path.exists(path_or_url):
                raise FileNotFoundError(f"Image not found: {path_or_url}")
            image = Image.open(path_or_url)

        text = pytesseract.image_to_string(image, lang='eng')
        return text.strip()

    except Exception as e:
        print(f"❌ Error processing image {path_or_url}: {e}")
        return ""


# print(extract_image("https://hackrx.blob.core.windows.net/assets/Test%20/image.png?sv=2023-01-03&spr=https&st=2025-08-04T19%3A21%3A45Z&se=2026-08-05T19%3A21%3A00Z&sr=b&sp=r&sig=lAn5WYGN%2BUAH7mBtlwGG4REw5EwYfsBtPrPuB0b18M4%3D"))
# print(extract_image("/home/krishna/hackathon-agent/files/image.png"))